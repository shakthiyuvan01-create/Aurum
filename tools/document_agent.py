"""
tools/document_agent.py — Document Intelligence Agent.
Upload large PDFs/DOCX/Excel. Summarize, Q&A, compare versions,
extract tables, create PowerPoint/Word report, generate action items.
"""
import json, logging, os, re
from pathlib import Path

log = logging.getLogger(__name__)

NAME        = "document_agent"
DESCRIPTION = (
    "Intelligent document processing. Actions: summarize, qa, compare, "
    "extract_tables, create_report, action_items, analyze. "
    "Supports PDF, DOCX, XLSX, PPTX."
)
CATEGORY = "documents"
ICON     = "📄"
INPUTS = [
    {"name": "action",    "label": "Action", "type": "select",
     "options": ["summarize","qa","compare","extract_tables","create_report","action_items","analyze"],
     "required": True},
    {"name": "file_path", "label": "File path or URL",   "type": "text"},
    {"name": "file_path2","label": "Second file (compare)", "type": "text"},
    {"name": "question",  "label": "Question (for Q&A)", "type": "text"},
    {"name": "output_format", "label": "Output format",  "type": "select",
     "options": ["text","docx","pptx","json"], "default": "text"},
    {"name": "username",  "label": "Username",           "type": "text"},
]


# ── Text extraction ────────────────────────────────────────────────────────────
def _extract_text(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    try:
        if ext == ".pdf":
            import pypdf
            reader = pypdf.PdfReader(file_path)
            return "\n".join(p.extract_text() or "" for p in reader.pages)

        if ext == ".docx":
            from docx import Document
            doc = Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs)

        if ext in (".xlsx", ".xls"):
            import openpyxl
            wb   = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            rows = []
            for ws in wb.worksheets:
                rows.append(f"[Sheet: {ws.title}]")
                for row in ws.iter_rows(values_only=True, max_row=200):
                    rows.append("\t".join(str(c or "") for c in row))
            return "\n".join(rows)

        if ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(file_path)
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [sh.text for sh in slide.shapes if hasattr(sh, "text") and sh.text.strip()]
                slides.append(f"Slide {i}: " + " | ".join(texts))
            return "\n".join(slides)

        if ext == ".txt":
            return open(file_path, encoding="utf-8", errors="ignore").read()

        return f"[Unsupported file type: {ext}]"
    except Exception as e:
        return f"[Extraction error: {e}]"


def _ai(prompt: str, max_tokens: int = 1500) -> str:
    from providers import AI
    return AI.generate(prompt, model="gpt-4o", max_tokens=max_tokens, temperature=0.3)


def _chunk_text(text: str, chunk_size: int = 8000) -> list[str]:
    words  = text.split()
    chunks = []
    buf    = []
    count  = 0
    for word in words:
        buf.append(word)
        count += 1
        if count >= chunk_size:
            chunks.append(" ".join(buf))
            buf   = []
            count = 0
    if buf:
        chunks.append(" ".join(buf))
    return chunks


# ── Actions ────────────────────────────────────────────────────────────────────
def _summarize(text: str, file_name: str) -> str:
    chunks = _chunk_text(text, 6000)
    if len(chunks) == 1:
        return _ai(f"Summarize this document comprehensively. File: {file_name}\n\n{chunks[0]}")
    # Map-reduce for large docs
    chunk_summaries = [
        _ai(f"Summarize this section (part {i+1}/{len(chunks)}):\n{chunk}", max_tokens=400)
        for i, chunk in enumerate(chunks[:6])
    ]
    combined = "\n\n".join(chunk_summaries)
    return _ai(f"Create a final comprehensive summary from these section summaries:\n\n{combined}", max_tokens=800)


def _qa(text: str, question: str, file_path: str = "", username: str = "") -> str:
    # RAG: chunk -> embed -> retrieve only the relevant parts (services/rag_service)
    try:
        from services.rag_service import qa_context
        relevant = qa_context(file_path or "inline", text, question,
                              username=username or "default")
    except Exception:
        relevant = _chunk_text(text, 6000)[0]
    return _ai(
        f"Answer this question based ONLY on the document text.\n"
        f"Question: {question}\n\nDocument:\n{relevant}",
        max_tokens=600,
    )


def _compare(text1: str, text2: str, name1: str, name2: str) -> str:
    t1 = text1[:3000]
    t2 = text2[:3000]
    return _ai(
        f"Compare these two documents.\n"
        f"Document 1 ({name1}):\n{t1}\n\n"
        f"Document 2 ({name2}):\n{t2}\n\n"
        "Provide: key similarities, key differences, what changed, recommendation.",
        max_tokens=800,
    )


def _extract_tables(text: str) -> str:
    return _ai(
        f"Extract ALL tables from this document text. Format each as a markdown table.\n\n{text[:8000]}",
        max_tokens=1200,
    )


def _action_items(text: str) -> str:
    return _ai(
        f"Extract ALL action items, tasks, decisions, and next steps from this document.\n"
        f"Format as numbered list with: Owner (if mentioned), Deadline (if mentioned), Action.\n\n{text[:8000]}",
        max_tokens=600,
    )


def _analyze(text: str, file_name: str) -> str:
    return _ai(
        f"Provide a comprehensive analysis of this document: {file_name}\n"
        "Include: Purpose, Key topics, Important data/facts, Sentiment, Quality assessment, Recommendations.\n\n"
        + text[:6000],
        max_tokens=1000,
    )


# ── Main ──────────────────────────────────────────────────────────────────────
def run(
    action:        str = "summarize",
    file_path:     str = "",
    file_path2:    str = "",
    question:      str = "",
    output_format: str = "text",
    username:      str = "",
) -> dict:
    action = (action or "summarize").lower().strip()

    if not file_path:
        return {"error": "file_path required"}

    if not os.path.exists(file_path):
        return {"error": f"File not found: {file_path}"}

    text  = _extract_text(file_path)
    fname = Path(file_path).name

    if action == "summarize":
        result = _summarize(text, fname)
        return {"result": result, "word_count": len(text.split())}

    if action == "qa":
        if not question:
            return {"error": "question required for Q&A"}
        result = _qa(text, question, file_path=file_path, username=username)
        return {"result": result}

    if action == "compare":
        if not file_path2 or not os.path.exists(file_path2):
            return {"error": "file_path2 required and must exist"}
        text2  = _extract_text(file_path2)
        fname2 = Path(file_path2).name
        result = _compare(text, text2, fname, fname2)
        return {"result": result}

    if action == "extract_tables":
        result = _extract_tables(text)
        return {"result": result}

    if action == "action_items":
        result = _action_items(text)
        return {"result": result}

    if action == "analyze":
        result = _analyze(text, fname)
        return {"result": result, "word_count": len(text.split())}

    if action == "create_report":
        summary      = _summarize(text, fname)
        action_items = _action_items(text)
        return {
            "result": f"## Document Report: {fname}\n\n### Summary\n{summary}\n\n### Action Items\n{action_items}",
            "summary":      summary,
            "action_items": action_items,
        }

    return {"error": f"Unknown action: {action}"}
