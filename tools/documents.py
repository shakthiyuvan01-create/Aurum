"""tools/documents.py - Document processing (Word, Excel, PPT, PDF, CSV, MD, HTML, OCR)"""
from __future__ import annotations
import csv, json, os
from pathlib import Path

NAME = "documents"
DESCRIPTION = "Read/create/edit: Word, Excel, PPT, PDF, CSV, Markdown, HTML, OCR"
CATEGORY = "builtin"
ICON = "📄"
INPUTS = [
    {"name": "action", "label": "Action", "type": "str", "required": True,
     "placeholder": "read_docx | create_docx | read_xlsx | create_xlsx | read_pptx | read_pdf | read_csv | create_csv | to_markdown | to_html | ocr_image | generate_report"},
    {"name": "path", "label": "File path", "type": "str"},
    {"name": "content", "label": "Content (JSON for create)", "type": "str"},
    {"name": "output_path", "label": "Output path", "type": "str"},
]

def run(action="", path="", content="", output_path=""):
    m = {"read_docx": _read_docx, "create_docx": _create_docx, "read_xlsx": _read_xlsx,
         "create_xlsx": _create_xlsx, "read_pptx": _read_pptx, "read_pdf": _read_pdf,
         "read_csv": _read_csv, "create_csv": _create_csv, "to_markdown": _to_markdown,
         "to_html": _to_html, "ocr_image": _ocr_image, "generate_report": _generate_report}
    fn = m.get(action.strip().lower())
    if fn: return fn(path=path, content=content, output_path=output_path)
    return {"error": f"Unknown: {action}"}

def _read_docx(path="", **kw):
    try:
        from docx import Document
        d = Document(path)
        return {"result": "\n".join(p.text for p in d.paragraphs), "paragraphs": len(d.paragraphs)}
    except ImportError: return {"error": "pip install python-docx"}
    except Exception as e: return {"error": str(e)}

def _create_docx(path="", content="", output_path="", **kw):
    try:
        from docx import Document
        d = Document()
        data = json.loads(content) if isinstance(content, str) else (content or {})
        if isinstance(data, list):
            for item in data:
                if isinstance(item, dict) and item.get("style")=="heading":
                    d.add_heading(item["text"], item.get("level",1))
                else: d.add_paragraph(str(item))
        elif isinstance(data, dict):
            d.add_heading(data.get("title","Doc"), 1)
            for p in data.get("body",[]): d.add_paragraph(p)
        out = output_path or "output.docx"
        d.save(out)
        return {"result": f"Created {out}", "file": out}
    except ImportError: return {"error": "pip install python-docx"}
    except Exception as e: return {"error": str(e)}

def _read_xlsx(path="", **kw):
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True)
        r = {}
        for s in wb.sheetnames:
            ws = wb[s]
            r[s] = [[str(c) if c is not None else "" for c in row] for row in ws.iter_rows(values_only=True)]
        wb.close()
        return {"result": r, "sheets": list(r.keys())}
    except ImportError: return {"error": "pip install openpyxl"}
    except Exception as e: return {"error": str(e)}

def _create_xlsx(path="", content="", output_path="", **kw):
    try:
        import openpyxl
        wb = openpyxl.Workbook()
        data = json.loads(content) if isinstance(content, str) else (content or {})
        if isinstance(data, dict):
            for i,(sn,rows) in enumerate(data.items()):
                ws = wb.active if i==0 else wb.create_sheet(title=sn); ws.title=sn
                for r in rows: ws.append(r if isinstance(r,list) else [r])
        elif isinstance(data, list):
            ws = wb.active
            for r in data: ws.append(r if isinstance(r,list) else [r])
        out = output_path or "output.xlsx"; wb.save(out); wb.close()
        return {"result": f"Created {out}", "file": out}
    except ImportError: return {"error": "pip install openpyxl"}
    except Exception as e: return {"error": str(e)}

def _read_pptx(path="", **kw):
    try:
        from pptx import Presentation
        prs = Presentation(path); slides=[]
        for s in prs.slides:
            t=[sh.text for sh in s.shapes if hasattr(sh,"text") and sh.text.strip()]
            slides.append("\n".join(t))
        return {"result": slides, "slide_count": len(slides)}
    except ImportError: return {"error": "pip install python-pptx"}
    except Exception as e: return {"error": str(e)}

def _read_pdf(path="", **kw):
    text_pages = []
    try:
        import PyPDF2
        with open(path,"rb") as f:
            for p in PyPDF2.PdfReader(f).pages:
                text_pages.append(p.extract_text() or "")
    except ImportError:
        try:
            import pdfplumber
            with pdfplumber.open(path) as pdf:
                for p in pdf.pages: text_pages.append(p.extract_text() or "")
        except ImportError: return {"error": "pip install PyPDF2 or pdfplumber"}
    except Exception as e: return {"error": str(e)}
    return {"result": "\n\n".join(text_pages), "pages": len(text_pages)}

def _read_csv(path="", **kw):
    try:
        rows=[]; f=open(path,newline="",encoding="utf-8")
        for r in csv.reader(f): rows.append(r)
        f.close()
        return {"result": rows, "row_count": len(rows)}
    except Exception as e: return {"error": str(e)}

def _create_csv(path="", content="", output_path="", **kw):
    try:
        data = json.loads(content) if isinstance(content, str) else (content or [])
        out = output_path or "output.csv"
        with open(out,"w",newline="",encoding="utf-8") as f:
            w=csv.writer(f)
            for r in (data if isinstance(data,list) else []):
                w.writerow(r if isinstance(r,list) else [r])
        return {"result": f"Created {out}", "file": out, "rows": len(data) if isinstance(data,list) else 0}
    except Exception as e: return {"error": str(e)}

def _to_markdown(path="", content="", output_path="", **kw):
    if output_path:
        with open(output_path,"w",encoding="utf-8") as f: f.write(content or "")
        return {"result": f"Written to {output_path}", "file": output_path}
    return {"result": content or ""}

def _to_html(path="", content="", output_path="", **kw):
    html = f"<html><body><pre>{content}</pre></body></html>"
    if output_path:
        with open(output_path,"w",encoding="utf-8") as f: f.write(html)
        return {"result": f"Written to {output_path}", "file": output_path}
    return {"result": html}

def _ocr_image(path="", **kw):
    try:
        import pytesseract; from PIL import Image
        return {"result": pytesseract.image_to_string(Image.open(path)).strip()}
    except ImportError: return {"error": "pip install pytesseract Pillow"}
    except Exception as e: return {"error": str(e)}

def _generate_report(path="", content="", output_path="", **kw):
    try: data = json.loads(content) if isinstance(content, str) else (content or {})
    except: data = {"body": [content]}
    title = data.get("title","Report"); sections = data.get("body",data.get("sections",[]))
    fmt = data.get("format","markdown")
    if fmt in ("markdown","md"):
        md = f"# {title}\n\n"+"\n".join(f"## {s.get('heading','')}\n\n{s.get('text','')}\n" if isinstance(s,dict) else f"{s}\n" for s in sections)
        out = output_path or "report.md"
        with open(out,"w",encoding="utf-8") as f: f.write(md)
        return {"result": f"Report: {out}", "file": out, "format":"markdown"}
    elif fmt=="html":
        html = f"<html><head><title>{title}</title></head><body><h1>{title}</h1>"
        for s in sections:
            html += f"<h2>{s.get('heading','')}</h2><p>{s.get('text','')}</p>" if isinstance(s,dict) else f"<p>{s}</p>"
        html += "</body></html>"
        out = output_path or "report.html"
        with open(out,"w",encoding="utf-8") as f: f.write(html)
        return {"result": f"Report: {out}", "file": out, "format":"html"}
    return {"error": f"Unknown format: {fmt}"}
